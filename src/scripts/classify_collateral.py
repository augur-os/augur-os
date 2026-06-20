#!/usr/bin/env python3
"""
Collateral Routing Script (ADR-135, Section 6)

Scans the repo root for stray files outside the canonical root layout,
builds a context bundle for the calling agent, and routes files to the
correct skill assets directories.

Two-step flow (the calling IDE agent IS the LLM):
    Step 1: python3 classify_collateral.py --mode context [--verbose]
            → stdout JSON: {stray_files, git_context, skill_registry, prompt}
    Step 2: Agent reads JSON, does inline classification
    Step 3: python3 classify_collateral.py --mode route --classification '{...}' [--dry-run] [--verbose]
            → Moves files to skill dirs, re-indexes, prints summary

Routing outcome:
    - Classified files  → skills/{skill}/assets/
    - _archive files    → external state garbage_collector/
"""

import argparse
import json
import shutil
import subprocess  # nosec B404
import sys
from pathlib import Path
from datetime import datetime

# ---------------------------------------------------------------------------
# Bootstrap sys.path so src.config.paths is importable when run directly
# ---------------------------------------------------------------------------
_SCRIPT_ROOT = Path(__file__).resolve().parents[2]  # project root
if str(_SCRIPT_ROOT) not in sys.path:
    sys.path.insert(0, str(_SCRIPT_ROOT))

from src.config.paths import get_project_brain_skills_dir, get_project_root, get_runtime_dir  # noqa: E402
from src.lib.repo_hygiene import collect_root_strays  # noqa: E402

# ---------------------------------------------------------------------------
# Text file extensions for content preview extraction
# ---------------------------------------------------------------------------
TEXT_EXTENSIONS = {".html", ".md", ".txt", ".csv", ".json", ".yaml", ".yml", ".py", ".ts", ".js"}


# ---------------------------------------------------------------------------
# File content extraction
# ---------------------------------------------------------------------------


def extract_text_preview(path: Path, max_chars: int = 500) -> str:
    """Return up to max_chars of text content from a file."""
    try:
        text = path.read_text(encoding="utf-8", errors="replace")
        return text[:max_chars]
    except Exception as e:
        return f"<read error: {e}>"


def extract_docx_preview(path: Path, max_chars: int = 500) -> str:
    """Extract text from a .docx file via python-docx if available."""
    try:
        import docx  # type: ignore

        doc = docx.Document(str(path))
        text = "\n".join(p.text for p in doc.paragraphs if p.text.strip())
        return text[:max_chars]
    except ImportError:
        return f"<python-docx not available; filename: {path.name}>"
    except Exception as e:
        return f"<docx read error: {e}>"


def extract_pptx_preview(path: Path, max_chars: int = 500) -> str:
    """Extract slide titles from a .pptx file via python-pptx if available."""
    try:
        from pptx import Presentation  # type: ignore

        prs = Presentation(str(path))
        titles = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame and shape.shape_type == 13:  # title placeholder
                    titles.append(shape.text_frame.text.strip())
                elif hasattr(shape, "text") and shape.text.strip():
                    titles.append(shape.text.strip()[:100])
                    break
        text = "; ".join(t for t in titles if t)
        return text[:max_chars] if text else f"<no text extracted; filename: {path.name}>"
    except ImportError:
        return f"<python-pptx not available; filename: {path.name}>"
    except Exception as e:
        return f"<pptx read error: {e}>"


def build_file_info(path: Path) -> dict:
    """Build a file descriptor dict for the classification prompt."""
    ext = path.suffix.lower()
    size_bytes = path.stat().st_size

    if ext in TEXT_EXTENSIONS:
        content_preview = extract_text_preview(path)
    elif ext == ".docx":
        content_preview = extract_docx_preview(path)
    elif ext == ".pptx":
        content_preview = extract_pptx_preview(path)
    else:
        content_preview = f"<binary; size: {size_bytes} bytes>"

    return {
        "filename": path.name,
        "extension": ext,
        "size_bytes": size_bytes,
        "content_preview": content_preview,
    }


# ---------------------------------------------------------------------------
# Git context
# ---------------------------------------------------------------------------


def get_git_context(project_root: Path) -> dict:
    """Return recent git log and diff stat for the classification prompt."""

    def run_git(args: list[str]) -> str:
        try:
            result = subprocess.run(  # nosec B603
                ["git"] + args,
                cwd=str(project_root),
                capture_output=True,
                text=True,
                timeout=10,
                encoding="utf-8",
            )
            return result.stdout.strip() if result.returncode == 0 else ""
        except Exception:
            return ""

    log_summary = run_git(["log", "--oneline", "-10"])
    diff_stat = run_git(["diff", "--stat", "HEAD~1"])
    branch_name = run_git(["branch", "--show-current"])

    return {
        "git_log_summary": log_summary or "<no commits>",
        "git_diff_stat": diff_stat or "<no diff>",
        "branch_name": branch_name or "main",
    }


# ---------------------------------------------------------------------------
# Skill registry
# ---------------------------------------------------------------------------


def get_skill_registry(project_root: Path) -> list[dict]:
    """
    Walk project-brain/capabilities/skills/ to collect canonical skill metadata.
    Returns list of {name, hub, description}.
    """
    skills = []
    skills_dir = get_project_brain_skills_dir(project_root)
    if not skills_dir.exists():
        return skills

    for skill_dir in sorted(skills_dir.iterdir()):
        if not skill_dir.is_dir() or skill_dir.name.startswith("."):
            continue
        skill_md = skill_dir / "SKILL.md"
        if not skill_md.exists():
            continue
        description = ""
        hub = "unknown"
        try:
            import yaml  # type: ignore

            text = skill_md.read_text(encoding="utf-8")
            if text.startswith("---"):
                end = text.index("---", 3)
                frontmatter = yaml.safe_load(text[3:end]) or {}
                if isinstance(frontmatter, dict):
                    description = str(frontmatter.get("description") or "")
            if not description:
                for line in text.splitlines()[:20]:
                    stripped = line.strip()
                    if stripped and not stripped.startswith("#") and not stripped.startswith("---"):
                        description = stripped[:120]
                        break
        except Exception:
            pass

        skills.append(
            {
                "name": skill_dir.name,
                "hub": hub,
                "description": description or f"Skill in {hub} hub",
            }
        )

    return skills


# ---------------------------------------------------------------------------
# Classification Prompt
# ---------------------------------------------------------------------------


def build_classification_prompt(
    stray_files: list[dict],
    git_context: dict,
    skill_registry: list[dict],
) -> str:
    """Assemble the classification prompt."""
    skill_lines = "\n".join(f"  - {s['name']} (hub: {s['hub']}): {s['description']}" for s in skill_registry)

    file_sections = []
    for f in stray_files:
        file_sections.append(
            f"### {f['filename']}\n"
            f"  Extension: {f['extension']}\n"
            f"  Size: {f['size_bytes']} bytes\n"
            f"  Content preview:\n  {f['content_preview'][:500]}"
        )
    files_block = "\n\n".join(file_sections)

    prompt = f"""Classify files into project skills. Return ONLY a JSON object.

Session context:
{git_context['git_log_summary']}

Branch: {git_context['branch_name']}

Skills (routing targets):
{skill_lines}

Files to classify:
{files_block}

Return a single JSON object where each key is a filename. Example:
{{
  "example.docx": {{"skill": "venture", "hub": "professional", "reason": "elevator pitch for venture"}},
  "temp.txt": {{"skill": "_archive", "hub": "_archive", "reason": "temporary scratch file"}}
}}

Rules:
- Keys must be exact filenames from the list above
- "skill" and "hub" must match a skill/hub from the list, or "_archive" for unclassifiable files
- Return ONLY the JSON object, no other text"""
    return prompt


# ---------------------------------------------------------------------------
# File routing
# ---------------------------------------------------------------------------


def route_files(
    classification: dict,
    stray_paths: dict[str, Path],
    project_root: Path,
    dry_run: bool = False,
    verbose: bool = False,
) -> dict:
    """
    Move each classified file to its target directory.

    Returns a summary dict with keys: routed, archived, errors.
    """
    summary: dict = {"routed": [], "archived": [], "errors": []}
    timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
    garbage_dir = get_runtime_dir() / "garbage_collector"

    for filename, info in classification.items():
        src_path = stray_paths.get(filename)
        if src_path is None or not src_path.exists():
            if verbose:
                print(f"[classify] Skip (already gone): {filename}", flush=True)
            continue

        skill = info.get("skill", "_archive")
        hub = info.get("hub", "_archive")
        reason = info.get("reason", "")

        if skill == "_archive":
            dest_dir = garbage_dir / timestamp
            label = f"{dest_dir}/"
        else:
            dest_dir = get_project_brain_skills_dir(project_root) / skill / "assets"
            label = f"project-brain/capabilities/skills/{skill}/assets/"

        dest_path = dest_dir / filename

        if dry_run:
            print(f"[DRY RUN] Would route: {filename} → {label}  ({reason})", flush=True)
            if skill == "_archive":
                summary["archived"].append({"file": filename, "dest": label, "reason": reason})
            else:
                summary["routed"].append(
                    {"file": filename, "dest": label, "skill": skill, "hub": hub, "reason": reason}
                )
        else:
            try:
                dest_dir.mkdir(parents=True, exist_ok=True)
                shutil.move(str(src_path), str(dest_path))
                if skill == "_archive":
                    print(f"[classify] Archived: {filename} → {label}", flush=True)
                    summary["archived"].append({"file": filename, "dest": label, "reason": reason})
                else:
                    print(f"[classify] Routed: {filename} → {label}  ({reason})", flush=True)
                    summary["routed"].append(
                        {"file": filename, "dest": label, "skill": skill, "hub": hub, "reason": reason}
                    )
            except Exception as e:
                print(f"[classify] ERROR routing {filename}: {e}", flush=True)
                summary["errors"].append({"file": filename, "error": str(e)})

    return summary


# ---------------------------------------------------------------------------
# Post-routing RAG re-index
# ---------------------------------------------------------------------------


def reindex_routed_skills(
    summary: dict,
    project_root: Path,
    dry_run: bool = False,
    verbose: bool = False,
) -> None:
    """
    Trigger RAG re-indexing for each skill that received routed files.

    Without this step, newly routed assets (especially binary files like
    .docx/.pptx that ripgrep cannot search inside) won't appear in the
    RAG index, causing the search to escalate to LLM unnecessarily.
    The unified_indexer generates category index entries that list filenames,
    making binary assets discoverable by ripgrep Tier 2 search.
    """
    if not summary.get("routed"):
        return

    # Collect unique skill directories that received files
    skill_dirs: set[str] = set()
    for entry in summary["routed"]:
        hub = entry.get("hub", "")
        skill = entry.get("skill", "")
        if hub and skill:
            skill_dirs.add(skill)

    if not skill_dirs:
        return

    if dry_run:
        for skill_rel in sorted(skill_dirs):
            print(f"[DRY RUN] Would re-index: skills/{skill_rel}", flush=True)
        return

    # Delegate to unified_indexer.reindex_all() — single pass covers all categories
    try:
        scripts_dir = get_project_brain_skills_dir(project_root) / "rag" / "scripts"
        if str(scripts_dir) not in sys.path:
            sys.path.insert(0, str(scripts_dir))

        from src.config.paths import get_documents_dir, get_rag_dir, get_vault_dir
        import unified_indexer

        if verbose:
            print(f"[classify] Running unified reindex for {len(skill_dirs)} routed skill(s)...", flush=True)
        stats = unified_indexer.reindex_all(
            root=project_root,
            rag_dir=get_rag_dir(),
            vault_dir=get_vault_dir(),
            documents_dir=get_documents_dir(),
        )
        total = sum(stats.values())
        if verbose:
            print(f"[classify] Unified reindex complete: {total} entries", flush=True)
    except Exception as e:
        print(f"[classify] WARNING: Unified reindex failed: {e}", flush=True)


# ---------------------------------------------------------------------------
# Mode: context — gather context and output JSON for the calling agent
# ---------------------------------------------------------------------------


def mode_context(project_root: Path, verbose: bool = False) -> None:
    """Gather stray files, git context, skill registry, and prompt. Output JSON to stdout."""
    # 1. Find stray files
    stray_paths: dict[str, Path] = {}
    for item in collect_root_strays(project_root):
        stray_paths[item.name] = item

    if not stray_paths:
        # Output empty context — no work to do
        output = {
            "stray_files": [],
            "git_context": {},
            "skill_registry": [],
            "prompt": "",
            "message": "No stray files found — repo root is clean.",
        }
        print(json.dumps(output, indent=2), flush=True)
        return

    if verbose:
        print(
            f"[classify] Found {len(stray_paths)} stray item(s): {', '.join(stray_paths.keys())}",
            file=sys.stderr,
            flush=True,
        )

    # 2. Build file info
    file_infos = []
    for name, path in stray_paths.items():
        try:
            info = build_file_info(path)
            file_infos.append(info)
        except Exception as e:
            if verbose:
                print(f"[classify] WARNING: Could not read {name}: {e}", file=sys.stderr, flush=True)
            file_infos.append(
                {
                    "filename": name,
                    "extension": Path(name).suffix.lower(),
                    "size_bytes": 0,
                    "content_preview": f"<unreadable: {e}>",
                }
            )

    # 3. Get git context and skill registry
    git_context = get_git_context(project_root)
    if verbose:
        print(f"[classify] Git branch: {git_context['branch_name']}", file=sys.stderr, flush=True)

    skill_registry = get_skill_registry(project_root)
    if verbose:
        print(f"[classify] Skill registry: {len(skill_registry)} skills found", file=sys.stderr, flush=True)

    # 4. Build prompt
    prompt = build_classification_prompt(file_infos, git_context, skill_registry)
    if verbose:
        print(f"[classify] Prompt length: {len(prompt)} chars", file=sys.stderr, flush=True)

    # 5. Output JSON to stdout (agent reads this)
    output = {
        "stray_files": file_infos,
        "git_context": git_context,
        "skill_registry": skill_registry,
        "prompt": prompt,
    }
    print(json.dumps(output, indent=2), flush=True)


# ---------------------------------------------------------------------------
# Mode: route — accept classification JSON and execute file routing
# ---------------------------------------------------------------------------


def mode_route(
    classification_json: str,
    project_root: Path,
    dry_run: bool = False,
    verbose: bool = False,
) -> None:
    """Parse classification JSON, route files, re-index, print summary."""
    # Parse classification
    try:
        classification = json.loads(classification_json)
    except json.JSONDecodeError as e:
        print(f"[classify] ERROR: Invalid classification JSON: {e}", flush=True)
        sys.exit(1)

    if not isinstance(classification, dict):
        print(
            f"[classify] ERROR: Classification must be a JSON object, got {type(classification).__name__}", flush=True
        )
        sys.exit(1)

    # Find stray files (need paths for routing)
    stray_paths: dict[str, Path] = {}
    for item in collect_root_strays(project_root):
        stray_paths[item.name] = item

    if verbose:
        print(f"[classify] Project root: {project_root}", flush=True)
        print(f"[classify] Mode: {'DRY RUN' if dry_run else 'LIVE'}", flush=True)
        print(f"[classify] Classification: {len(classification)} file(s)", flush=True)
        print(f"[classify] Stray files on disk: {len(stray_paths)}", flush=True)

    # Route files
    summary = route_files(classification, stray_paths, project_root, dry_run, verbose)

    # Re-index skills that received new files
    reindex_routed_skills(summary, project_root, dry_run, verbose)

    # Print summary
    print("-" * 50, flush=True)
    print(f"[classify] Routed:   {len(summary['routed'])} file(s)", flush=True)
    print(f"[classify] Archived: {len(summary['archived'])} file(s)", flush=True)
    print(f"[classify] Errors:   {len(summary['errors'])} file(s)", flush=True)

    if summary["routed"]:
        print("\nRouted files (commit with: chore(assets): route session collateral to skill data dirs):", flush=True)
        for entry in summary["routed"]:
            print(f"  {entry['file']} → {entry['dest']}", flush=True)

    if summary["errors"]:
        sys.exit(1)


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------


def main() -> None:
    parser = argparse.ArgumentParser(
        description="Collateral routing — classify and route stray root files to skill assets directories."
    )
    parser.add_argument(
        "--mode",
        choices=["context", "route"],
        default="context",
        help="Operation mode: 'context' gathers stray files and outputs JSON for agent classification; "
        "'route' accepts a classification dict and executes file routing (default: context)",
    )
    parser.add_argument(
        "--classification",
        type=str,
        default=None,
        help="JSON classification dict (required for --mode route). "
        'Example: \'{"file.docx": {"skill": "venture", "hub": "professional", "reason": "..."}}\'',
    )
    parser.add_argument("--dry-run", action="store_true", help="Simulate routing without moving files")
    parser.add_argument(
        "--root-dir", type=str, default=None, help="Project root directory (defaults to auto-detected project root)"
    )
    parser.add_argument("--verbose", action="store_true", help="Print detailed progress to stderr")
    args = parser.parse_args()

    project_root = Path(args.root_dir).resolve() if args.root_dir else get_project_root()

    if args.mode == "context":
        mode_context(project_root, args.verbose)

    elif args.mode == "route":
        if not args.classification:
            parser.error("--classification is required for --mode route")
        mode_route(args.classification, project_root, args.dry_run, args.verbose)


if __name__ == "__main__":
    main()
